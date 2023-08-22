import comtypes.client
import sys
import csv
import collections
import collections.abc
from pptx import Presentation
import os
import smtplib
from email.message import EmailMessage
import glob
import os
import tqdm
import time
import os
import sys
import comtypes.client


company_name = ""
email1 = ""
email2 = ""
email3 = ""
email4 = ""

with open('data.csv', 'r') as csv_file:
    csv_reader = csv.reader(csv_file)
    next(csv_reader)
    for index, line in enumerate(csv_reader):
        print("Index:", index)
        company_name = line[7]
        email1 = line[16]
        email2 = line[19]
        email3 = line[22]
        email4 = line[25]
        if (company_name == ""):
            continue

    # PPT Change Text
        prs = Presentation('test.pptx')

        # To get shapes in your slides
        slides = [slide for slide in prs.slides]
        shapes = []
        for slide in slides:
            for shape in slide.shapes:
                shapes.append(shape)

        def replace_text(replacements: dict, shapes: list, company):
            """Takes dict of {match: replacement, ... } and replaces all matches.
            Currently not implemented for charts or graphics.
            """
            for shape in shapes:
                for match, replacement in replacements.items():
                    if shape.has_text_frame:
                        if (shape.text.find(match)) != -1:
                            text_frame = shape.text_frame
                            for paragraph in text_frame.paragraphs:
                                for run in paragraph.runs:
                                    cur_text = run.text
                                    new_text = cur_text.replace(
                                        str(match), str(replacement))
                                    run.text = new_text
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if match in cell.text:
                                    new_text = cell.text.replace(
                                        match, replacement)
                                    cell.text = new_text

            prs.save(f'proposals/{company}.pptx')

        replace_text({'Durian': company_name}, shapes, company_name)

        print("PPT Done")

        # Email

        def PPTtoPDF(inputFileName, outputFileName, formatType=32):
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = 1

            if outputFileName[-3:] != 'pdf':
                outputFileName = outputFileName + ".pdf"
            deck = powerpoint.Presentations.Open(inputFileName)
            # formatType = 32 for ppt to pdf
            deck.SaveAs(outputFileName, formatType)
            deck.Close()
            powerpoint.Quit()

        # Open the ppt file path. sys.argv[1] is the ppt name.
            mypath = os.path.abspath(_file_)
            mydir = os.path.dirname(mypath)
            file_input = os.path.join(mydir, f'proposals/{company}.pptx')

        # create the pdf output file path and call your function
            file_output = os.path.join(
                mydir, f'proposals/{company}.pptx'[1][:-4] + "pdf")
            PPTtoPDF(file_input, file_output)

            time.sleep(1)

        EMAIL_ADDRESS = "ayushgoyal.thomso@gmail.com"
        EMAIL_PASSWORD = "your-key"

        contacts = [email1, email2, email3, email4]

        msg = EmailMessage()
        msg['Subject'] = f'Proposal for association between {company_name} and Thomso, IIT Roorkee '
        msg['From'] = EMAIL_ADDRESS
        msg['To'] = contacts
        msg['Cc'] = "chinmay.thomso@gmail.com"

        msg.set_content(f''' Dear Sir/Madam,
            Greetings from Thomso, IIT Roorkee!
            We are honoured to invite {company_name} as a Sponsor for Thomso '23, the annual cultural fest of IIT Roorkee. It is our marquee event, which attracts over 1,00,000 people every year. This time around it’s going to be bigger and better than ever before! We have been proud hosts to great personalities like Mrs. Smriti Irani (as our Chief Guest), Sonu Nigam, Darshan Rawal, Farhan Akhtar, Salim-Sulaiman, and Jubin Nautiyal (check it out: https://youtu.be/h7gyJRWrjbg )in our previous editions (Thomso ’22: https://youtu.be/rm1bWDAHbSQ ). Kindly check the attached brochure for more information about Thomso.
            With several zonal events lined up at premium locations like Delhi, Lucknow, Jaipur, Chandigarh, Bangalore, and Ahmedabad, Thomso's reach is more than 2,00,000 people through various Newspapers, Magazines, Social Media Platforms, Radio and T.V. Channels, Online Blogs, etc.
            This is a unique opportunity for {company_name} to engage, to introduce itself to prospective customers as well as promote growth strategies with existing clientele. The attached proposal  contains details regarding sponsorship prospects, deliverables, etc.
            All the Deliverables and related things are negotiable from both parties and can be discussed further on mail or call.
            We have come a long way but would still greatly benefit from your support. With immense participation and massive outreach, this affiliation will prove to be mutually beneficial.

            Inviting you to be a part of our "THOMSO" family of IIT ROORKEE.
            ''')

        files = ['Brochure.pdf', f'proposals/{company_name}.pdf']
        for file in files:
            with open(file, 'rb') as f:
                file_data = f.read()
                file_name = f.name

            msg.add_attachment(file_data, maintype='application',
                               subtype='octet-stream', filename=file_name)
            msg.add_header('X-Unsent', '1')
            # msg.add_alternative("""\
            # <!DOCTYPE html>
            # <html>
            #     <body>
            #         <h1 style="color:SlateGray;">This is an HTML Email!</h1>
            #     </body>
            # </html>
            # """, subtype='html')

            print("attachment added")
            with smtplib.SMTP_SSL('smtp.gmail.com', 465) as smtp:
                smtp.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
                print("Logged In")
                smtp.send_message(msg)

            print("Email Sent")
